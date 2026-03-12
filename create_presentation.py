from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

# Define slide contents
slides_data = [
    {
        "title": "Natural Acceptance and Harmonious Coexistence: A Case Study of Dujiangyan",
        "subtitle": "Exploring Mutually Fulfilling Relationships Across the Four Orders of Nature\n\nPresenter Name: arka2024\nDate: March 12, 2026\n\n[Insert High-Res Aerial Photo of Dujiangyan here]"
    },
    {
        "title": "Deconstructing Nature: The Four Orders of Coexistence",
        "content": "The Foundation: To analyze harmonious coexistence, we must define the four interdependent orders of nature:\n"
                   "1. Material Order: Inanimate foundation (soil, water, rocks). Characterized by composition/decomposition.\n"
                   "2. Plant Order: The biosystem (forests, crops). Characterized by growth and respiration.\n"
                   "3. Animal Order: Mobile life (fish, birds). Characterized by consciousness and physical activity.\n"
                   "4. Human Order: The realizing order. Characterized by rational thought, conscious choice, and 'natural acceptance.'\n\n"
                   "The Core Argument: Human fulfillment is achieved not through exploitation, but by ensuring our activities are mutually fulfilling to the other three orders.\n\n"
                   "[Insert Circular Diagram of the 4 Orders here]"
    },
    {
        "title": "The Challenge: Taming the Wild Minjiang River",
        "content": "The Context: Built in 256 BCE in the Sichuan Basin, China, by Governor Li Bing.\n\n"
                   "The Problem: The steep descent of the Minjiang River caused massive spring flash floods, destroying settlements (Human Order) and crops (Plant Order), while the chaotic water/debris (Material Order) wreaked havoc.\n\n"
                   "The Paradigm Shift: Instead of fighting the river with brute force (traditional dams/walls), Li Bing developed a hydraulic intervention that used the river's own energy to regulate itself.\n\n"
                   "The Goal: To examine a 2,200-year-old system where mutual fulfillment is an engineered reality, transforming a disaster zone into a 'Land of Abundance.'\n\n"
                   "[Insert Topographical Map of Minjiang River here]"
    },
    {
        "title": "Technical Marvels of Coexistence",
        "content": "Harmonizing the Material Order (The Fish Mouth): A tear-drop shaped man-made island that splits the river into an inner channel (irrigation) and outer channel (flood/sediment discharge). It uses centrifugal force to naturally flush heavy silt without blocking the river.\n\n"
                   "Harmonizing the Plant Order (The Flying Sand Weir): A spillway that uses vortex dynamics to eject excess sediment. This delivers clean water to the plains, preventing soil salinization and allowing crops to thrive.\n\n"
                   "Harmonizing the Animal Order (The Bottle-Neck Channel): An open-channel flow valve carved through the mountain. Unlike modern dams that block aquatic migration, this open system allows fish and aquatic life to move freely, preserving biodiversity.\n\n"
                   "[Insert Technical Diagram of Water Flow here]"
    },
    {
        "title": "Governance Rooted in 'Following Nature' (Dao)",
        "content": "The Philosophy of 'Wu Wei': The builders followed the Taoist principle of Dao Fa Zi Ran (Nature follows the Dao). They practiced strategic non-intervention—observing the natural flow and providing minimal, intelligent guidance.\n\n"
                   "Natural Acceptance in Practice: The engineers viewed the river as a partner, not an enemy. Technology was restricted by environmental ethics; the tool served life rather than dominating it.\n\n"
                   "The 'Sacred Rule' of Maintenance: Carved into the riverbed: 'Dig the bed deep, keep the dykes low.' This mandates annual dredging to help the river flow naturally, rather than building higher walls to trap it.\n\n"
                   "[Insert Historical Painting of Dredging Ritual here]"
    },
    {
        "title": "Synergy Over 2,000 Years",
        "content": "The 2,200-Year Stress Test: Dujiangyan has operated continuously without failing, proving that nature is self-regulating when humans participate intelligently.\n\n"
                   "Modern Engineering vs. Ancient Wisdom:\n"
                   "- Modern Dams: Often view rivers as obstacles. They block the Material Order, starve deltas of sediment, and fragment ecosystems.\n"
                   "- Dujiangyan: Views the river as a flow. It guides the Material Order, enriches the Plant Order, and protects the corridors of the Animal Order.\n\n"
                   "The Ultimate Lesson: Fulfilling human needs does not require the destruction of nature. True technological advancement listens to nature rather than conquering it.\n\n"
                   "[Insert Split-Screen: Modern Dam vs. Dujiangyan here]"
    },
    {
        "title": "Conclusion & References",
        "content": "Final Thought: Dujiangyan is a living testament that there is a provision in nature for the harmonious coexistence of all four orders. The Human Order's highest function is to understand, nurture, and participate in this larger harmony.\n\n"
                   "Key References:\n"
                   "- Huang, B., et al. (2024). Governing human-nature harmony in protected area communities. Bulletin of Chinese Academy of Sciences.\n"
                   "- e-Adhyayan. (n.d.). Module on Concord in Nature: Understanding the Four Orders.\n"
                   "- Springer. (2019). Synergetic Evolution of Social and Natural Systems.\n"
                   "- Fu, H. (2020). The Innovative Road to Urban Nature Conservation.\n\n"
                   "[Insert Sunset View of Dujiangyan here]"
    }
]

# Add Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = slides_data[0]["title"]
subtitle.text = slides_data[0]["subtitle"]

# Add Content Slides
bullet_slide_layout = prs.slide_layouts[1]
for slide_data in slides_data[1:]:
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = slide_data["title"]
    
    tf = body_shape.text_frame
    tf.text = slide_data["content"]
    # Adjust font size to fit
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(16)

# Save the presentation
filename = "Dujiangyan_Presentation.pptx"
prs.save(filename)
print(f"Presentation successfully saved as {filename}")
